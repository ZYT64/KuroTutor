"""通用文档能力：Agent 的 Word / PPT / PDF 读写编辑服务。

覆盖教学场景的文档操作：
- **读**：docx（段落+样式）、pptx（逐页文字）、pdf（逐页文本）→ 结构化文本，供 Agent 理解内容；
- **写**：轻量标记（`# ` 标题 / `## ` 节·幻灯片 / `- ` 列表 / 普通段落）→ docx / pptx / pdf；
- **编辑**：docx 追加与替换文本、pptx 追加幻灯片、pdf 合并与抽页；
- **转换**：其余格式经 LibreOffice 无头模式互转。

设计红线：与 layout/llm 一致——可插拔、错误信息含「现象+原因+建议」、不硬编码路径。
"""

from __future__ import annotations

import re
from pathlib import Path

from kurotutor.core.errors import ProviderError

_PDF_FONT = "china-s"  # PyMuPDF 内置中文字体


def _exists(path: str) -> Path:
    p = Path(path)
    if not p.exists():
        raise ProviderError("文件不存在", cause=path, fix="确认路径（相对工作区或绝对路径）")
    return p


# ---- 读 ----------------------------------------------------------------------


def read_document(path: str, *, max_chars: int = 20000) -> str:
    """读文档为结构化文本。支持 .docx / .pptx / .pdf / .txt / .md。"""
    p = _exists(path)
    suffix = p.suffix.lower()
    if suffix == ".docx":
        return _read_docx(p, max_chars)
    if suffix == ".pptx":
        return _read_pptx(p, max_chars)
    if suffix == ".pdf":
        return _read_pdf(p, max_chars)
    if suffix in (".txt", ".md"):
        return p.read_text(encoding="utf-8", errors="replace")[:max_chars]
    raise ProviderError(
        f"暂不支持直接读取 {suffix} 格式", cause=p.name, fix="先用 doc_convert 转为 docx/pdf/pptx 再读"
    )


def _read_docx(p: Path, max_chars: int) -> str:
    import docx

    d = docx.Document(str(p))
    parts: list[str] = []
    used = 0
    for para in d.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = (para.style.name or "").lower()
        if "heading" in style or "title" in style:
            parts.append(f"[标题] {t}")
        else:
            parts.append(t)
        used += len(t) + 4
        if used >= max_chars:
            parts.append("…（已达读取上限，内容有截断）")
            break
    # 表格也读出来（试卷/讲义常见）
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts) or "（文档为空）"


def _read_pptx(p: Path, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts: list[str] = []
    used = 0
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"===== 第 {i} 页 =====")
        used += 12
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    parts.append(t)
                    used += len(t) + 1
        if used >= max_chars:
            parts.append("…（已达读取上限）")
            break
    return "\n".join(parts) or "（演示文稿为空）"


def _read_pdf(p: Path, max_chars: int) -> str:
    import pymupdf

    doc = pymupdf.open(str(p))
    parts: list[str] = []
    used = 0
    try:
        for i, page in enumerate(doc, 1):
            parts.append(f"===== 第 {i} 页 =====")
            text = page.get_text("text").strip()
            parts.append(text or "(本页无文本层，是扫描件——请改用 ocr_read 工具识别本页文字)")
            used += len(text) + 16
            if used >= max_chars:
                parts.append("…（已达读取上限）")
                break
    finally:
        doc.close()
    return "\n".join(parts) or "（PDF 为空）"


# ---- 写（轻量标记 → 文档） ----------------------------------------------------
# 标记规则：`# ` 大标题；`## ` 节标题（pptx 中为新一页幻灯片标题）；`- ` 列表项；普通行 = 段落。


def write_document(out_path: str, content: str) -> str:
    """把轻量标记内容写成文档。格式由 out_path 后缀决定（.docx / .pptx / .pdf）。"""
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    suffix = out.suffix.lower()
    blocks = _parse_markup(content)
    if suffix == ".docx":
        return _write_docx(out, blocks)
    if suffix == ".pptx":
        return _write_pptx(out, blocks)
    if suffix == ".pdf":
        return _write_pdf(out, blocks)
    raise ProviderError(
        f"暂不支持生成 {suffix} 格式",
        cause=out.name,
        fix="用 .docx / .pptx / .pdf 后缀，或先写 docx 再 doc_convert",
    )


def _parse_markup(content: str) -> list[tuple[str, str]]:
    """解析轻量标记 → [(kind, text)]，kind ∈ title / section / bullet / para / image。"""
    blocks: list[tuple[str, str]] = []
    for raw in (content or "").splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        m = re.match(r"^!\[[^\]]*\]\(([^)]+)\)\s*$", line.strip())
        if m:
            blocks.append(("image", m.group(1).strip()))
        elif line.startswith("# "):
            blocks.append(("title", line[2:].strip()))
        elif line.startswith("## "):
            blocks.append(("section", line[3:].strip()))
        elif line.startswith("- ") or line.startswith("• "):
            blocks.append(("bullet", line[2:].strip()))
        else:
            blocks.append(("para", line.strip()))
    return blocks


def _write_docx(out: Path, blocks: list[tuple[str, str]]) -> str:
    import docx
    from docx.shared import Inches

    d = docx.Document()
    for kind, text in blocks:
        if kind == "title":
            d.add_heading(text, level=0 if text == blocks[0][1] else 1)
        elif kind == "section":
            d.add_heading(text, level=2)
        elif kind == "bullet":
            d.add_paragraph(text, style="List Bullet")
        elif kind == "image":
            if Path(text).exists():
                d.add_picture(text, width=Inches(5.8))
        else:
            d.add_paragraph(text)
    d.save(str(out))
    return str(out)


def _write_pptx(out: Path, blocks: list[tuple[str, str]]) -> str:
    from pptx import Presentation
    from pptx.util import Pt as PPt

    prs = Presentation()
    slide = None
    body = None

    def _new_slide(title: str):
        nonlocal slide, body
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.word_wrap = True

    first_title = next((t for k, t in blocks if k == "title"), "演示文稿")
    _new_slide(first_title)
    for kind, text in blocks:
        if kind == "title" and text != first_title:
            _new_slide(text)
            continue
        if kind == "section":
            if body.paragraphs and (body.text or "").strip():
                _new_slide(text)
            else:
                slide.shapes.title.text = text
            continue
        last = body.paragraphs[-1]
        para = last if not last.text and not last.runs else body.add_paragraph()
        para.text = ("• " + text) if kind == "bullet" else text
        para.font.size = PPt(16)
    prs.save(str(out))
    return str(out)


def _write_pdf(out: Path, blocks: list[tuple[str, str]]) -> str:
    import pymupdf

    A4_W, A4_H, MARGIN, BOTTOM = 595, 842, 50, 60
    doc = pymupdf.open()
    page = doc.new_page(width=A4_W, height=A4_H)
    y = MARGIN

    def _ensure(h: float) -> None:
        nonlocal page, y
        if y + h > A4_H - BOTTOM:
            page = doc.new_page(width=A4_W, height=A4_H)
            y = MARGIN

    def _wrap(text: str, fontsize: float) -> list[str]:
        per = max(1, int((A4_W - 2 * MARGIN) / fontsize))
        return re.findall(rf".{{1,{per}}}", text)

    for kind, text in blocks:
        if kind == "title":
            _ensure(34)
            page.insert_text((MARGIN, y + 14), text, fontname=_PDF_FONT, fontsize=16)
            y += 30
        elif kind == "image":
            img = Path(text)
            if not img.exists():
                continue
            import pymupdf as _pm

            with _pm.open(str(img)) as idoc:
                iw, ih = idoc[0].rect.width, idoc[0].rect.height
            max_w = A4_W - 2 * MARGIN
            scale = min(max_w / iw, 1.0)
            w, h = iw * scale, ih * scale
            _ensure(h + 12)
            page.insert_image(_pm.Rect(MARGIN, y, MARGIN + w, y + h), filename=str(img))
            y += h + 12
        elif kind == "section":
            _ensure(28)
            page.insert_text((MARGIN, y + 12), text, fontname=_PDF_FONT, fontsize=13)
            y += 24
        elif kind == "bullet":
            for i, chunk in enumerate(_wrap(text, 11)):
                _ensure(18)
                prefix = "• " if i == 0 else "  "
                x = MARGIN + (0 if i == 0 else 14)
                page.insert_text((x, y), prefix + chunk, fontname=_PDF_FONT, fontsize=11)
                y += 16
        else:
            for chunk in _wrap(text, 11):
                _ensure(18)
                page.insert_text((MARGIN, y), chunk, fontname=_PDF_FONT, fontsize=11)
                y += 16
        y += 4
    doc.save(str(out))
    doc.close()
    return str(out)


# ---- 编辑 ---------------------------------------------------------------------


def edit_document(path: str, op: str, payload: str, *, replacement: str = "") -> str:
    """编辑既有文档。op：
    - ``append``：把轻量标记内容追加到 docx/pptx 末尾；
    - ``replace``：docx 内把 payload（原文）替换为 replacement（全文匹配段落级）。
    """
    p = _exists(path)
    suffix = p.suffix.lower()
    if op == "append" and suffix == ".docx":
        import docx

        d = docx.Document(str(p))
        for kind, text in _parse_markup(payload):
            if kind == "title":
                d.add_heading(text, level=1)
            elif kind == "section":
                d.add_heading(text, level=2)
            elif kind == "bullet":
                d.add_paragraph(text, style="List Bullet")
            else:
                d.add_paragraph(text)
        d.save(str(p))
        return str(p)
    if op == "append" and suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(p))
        blocks = _parse_markup(payload)
        for kind, text in blocks:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = text[:60] if kind in ("title", "section") else "新内容"
            body = slide.placeholders[1].text_frame
            body.text = text
        prs.save(str(p))
        return str(p)
    if op == "replace" and suffix == ".docx":
        import docx

        d = docx.Document(str(p))
        hits = 0
        for para in d.paragraphs:
            if payload and payload in para.text:
                for run in para.runs:
                    if payload in run.text:
                        run.text = run.text.replace(payload, replacement)
                        hits += 1
        if hits == 0:
            raise ProviderError(
                "没有找到要替换的文本", cause=f"未在 {p.name} 段落中匹配到：{payload[:40]}",
                fix="先用 doc_read 查看原文，按完整短语替换",
            )
        d.save(str(p))
        return f"{p}（替换 {hits} 处）"
    raise ProviderError(
        "该操作不支持此格式", cause=f"op={op} × {suffix}", fix="append 支持 docx/pptx；replace 仅支持 docx"
    )


def pdf_merge(paths: list[str], out_path: str) -> str:
    """合并多个 PDF（按顺序）。"""
    import pymupdf

    if len(paths) < 2:
        raise ProviderError("合并至少需要 2 个文件", cause=str(paths), fix="提供 2 个以上 PDF 路径")
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    merged = pymupdf.open()
    for p in paths:
        merged.insert_pdf(pymupdf.open(str(_exists(p))))
    merged.save(str(out))
    merged.close()
    return str(out)


def pdf_extract_pages(path: str, pages: str, out_path: str) -> str:
    """抽取 PDF 指定页（如 ``1,3,5-8``，1 起）另存。"""
    import pymupdf

    src = pymupdf.open(str(_exists(path)))
    picked: list[int] = []
    for part in re.split(r"[,，]", pages or ""):
        part = part.strip()
        if not part:
            continue
        m = re.match(r"^(\d+)\s*-\s*(\d+)$", part)
        if m:
            picked.extend(range(int(m.group(1)) - 1, int(m.group(2))))
        elif part.isdigit():
            picked.append(int(part) - 1)
    total = src.page_count
    picked = [i for i in picked if 0 <= i < total]
    if not picked:
        src.close()
        raise ProviderError(
            "没有有效的页码", cause=f"文档共 {total} 页，解析结果为空", fix="用 1,3,5-8 这样的页码格式"
        )
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    dst = pymupdf.open()
    for i in picked:
        dst.insert_pdf(src, from_page=i, to_page=i)
    dst.save(str(out))
    dst.close()
    src.close()
    return f"{out}（共 {len(picked)} 页）"


def convert_document(path: str, target_dir: str, target_fmt: str) -> str:
    """经 LibreOffice 无头模式转换格式（docx/pptx/pdf/txt 等互转）。返回新文件路径。"""
    import shutil as _shutil
    import subprocess

    p = _exists(path)
    soffice = _shutil.which("soffice") or _shutil.which("soffice.exe")
    if not soffice:
        raise ProviderError(
            "格式转换需要 LibreOffice",
            cause="未在 PATH 中找到 soffice",
            fix="安装 LibreOffice（apt install libreoffice / 官网下载），或改用 doc_write 直接生成目标格式",
        )
    fmt = target_fmt.lstrip(".").lower()
    out = Path(target_dir)
    out.mkdir(parents=True, exist_ok=True)
    proc = subprocess.run(
        [soffice, "--headless", "--convert-to", fmt, "--outdir", str(out), str(p)],
        capture_output=True,
        text=True,
        timeout=180,
    )
    result = out / (p.stem + "." + fmt)
    if proc.returncode != 0 or not result.exists():
        raise ProviderError(
            "LibreOffice 转换失败",
            cause=(proc.stderr or proc.stdout)[:200],
            fix="确认目标格式受支持且文档未加密",
        )
    return str(result)
